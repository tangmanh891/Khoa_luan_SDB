"""Sinh slide thuyết trình cho khóa luận AutoShotV2.

Chạy:
    python publications/thesis/slides/build_slides.py

Output:
    publications/thesis/build/AutoShotV2_Defense.pptx
    Có thể đổi bằng tham số --output.

Yêu cầu:
    pip install python-pptx

Ghi chú:
    - Khổ 16:9 (13.333 x 7.5 inch).
    - Số liệu lấy từ publications/thesis/generated/slide_results.json.
    - Hình nền từ ../images/.
    - Font ưu tiên Calibri (có sẵn Windows, hỗ trợ tiếng Việt đầy đủ).
"""
from __future__ import annotations

import argparse
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "..", "images")
RESULTS_PATH = os.path.join(HERE, "..", "generated", "slide_results.json")

with open(RESULTS_PATH, "r", encoding="utf-8") as handle:
    RESULTS = json.load(handle)

COMPARISONS = {item["id"]: item for item in RESULTS["comparison_models"]}
ABLATIONS = {item["id"]: item for item in RESULTS["ablation"]}


def comparison_f1(identifier, dataset):
    return COMPARISONS[identifier]["metrics"][dataset]


def ablation_f1(identifier, dataset):
    return ABLATIONS[identifier]["metrics"][dataset]["f1"]


def f1_text(value):
    return "—" if value is None else f"{value:.3f}"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLOR_BG      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PANEL   = RGBColor(0xF4, 0xF6, 0xF8)
COLOR_ACCENT  = RGBColor(0x0B, 0x5F, 0xFF)
COLOR_ACCENT2 = RGBColor(0xFF, 0x7A, 0x45)
COLOR_INK     = RGBColor(0x1F, 0x2A, 0x44)
COLOR_SUB     = RGBColor(0x5B, 0x6B, 0x85)
COLOR_LINE    = RGBColor(0xE1, 0xE6, 0xEE)
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT_MAIN  = "Calibri"
FONT_TITLE = "Calibri"

TOTAL_SLIDES = 27
FOOTER_TEXT  = "AutoShotV2  |  HCMUS 04/2026"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color: RGBColor, width_pt: float = 0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_text(slide, left, top, width, height, text, *,
             font=FONT_MAIN, size=14, bold=False, color=COLOR_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=COLOR_INK, bullet_color=COLOR_ACCENT,
                line_spacing=1.25):
    """Add a list of bullet points. Each item: str or (bold_prefix, rest)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # bullet
        b = p.add_run()
        b.text = "•  "
        b.font.name = FONT_MAIN
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name = FONT_MAIN
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = tail
            r2.font.name = FONT_MAIN
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = FONT_MAIN
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, *, fill=None, line=None,
             line_width=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.shadow.inherit = False
    if fill is not None:
        set_fill(sh, fill)
    else:
        sh.fill.background()
    if line is None:
        no_line(sh)
    else:
        set_line(sh, line, line_width)
    sh.text_frame.margin_left = Emu(0)
    sh.text_frame.margin_right = Emu(0)
    return sh


def add_round_rect(slide, left, top, width, height, *, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.adjustments[0] = 0.12
    sh.shadow.inherit = False
    if fill is not None:
        set_fill(sh, fill)
    else:
        sh.fill.background()
    if line is None:
        no_line(sh)
    else:
        set_line(sh, line)
    return sh


def add_line(slide, x1, y1, x2, y2, color=COLOR_LINE, width=1.0):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def add_image_fit(slide, path, left, top, width, height):
    """Add image centered inside the bounding box, preserving aspect ratio."""
    # python-pptx auto-scales if we pass both width and height, but that
    # distorts aspect ratio. Compute ratio manually via a temporary add.
    pic = slide.shapes.add_picture(path, left, top, width=width)
    # If height overshoots, rescale.
    if pic.height > height:
        ratio = height / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = height
    # Center inside the box.
    pic.left = int(left + (width - pic.width) / 2)
    pic.top = int(top + (height - pic.height) / 2)
    return pic


# ---------------------------------------------------------------------------
# Slide chrome (title, footer, accent)
# ---------------------------------------------------------------------------

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_BG)
    return slide


def add_chrome(slide, title: str, kicker: str | None, page_num: int):
    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=COLOR_ACCENT)
    # Kicker (small label above title)
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.35), Inches(10), Inches(0.3),
                 kicker.upper(), size=11, bold=True, color=COLOR_ACCENT,
                 font=FONT_TITLE)
    # Title
    add_text(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=COLOR_INK, font=FONT_TITLE)
    # Thin rule under title
    add_line(slide, Inches(0.6), Inches(1.45), Inches(12.73), Inches(1.45),
             color=COLOR_LINE, width=0.75)
    # Footer rule
    add_line(slide, Inches(0.6), Inches(7.15), Inches(12.73), Inches(7.15),
             color=COLOR_LINE, width=0.5)
    # Footer text (left)
    add_text(slide, Inches(0.6), Inches(7.2), Inches(8), Inches(0.3),
             FOOTER_TEXT, size=9, color=COLOR_SUB)
    # Page number (right)
    add_text(slide, Inches(11.5), Inches(7.2), Inches(1.23), Inches(0.3),
             f"{page_num} / {TOTAL_SLIDES}", size=9, color=COLOR_SUB,
             align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_cover(prs):
    slide = blank_slide(prs)
    # Full colored left band
    add_rect(slide, 0, 0, Inches(4.5), SLIDE_H, fill=COLOR_ACCENT)
    # Decorative accent square
    add_rect(slide, Inches(3.7), Inches(6.2), Inches(0.8), Inches(0.8),
             fill=COLOR_ACCENT2)

    # Left band text
    add_text(slide, Inches(0.6), Inches(0.7), Inches(3.8), Inches(0.4),
             "KHÓA LUẬN TỐT NGHIỆP", size=13, bold=True,
             color=COLOR_WHITE, font=FONT_TITLE)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(3.8), Inches(0.4),
             "HCMUS  |  Khoa CNTT", size=11, color=COLOR_WHITE)

    # Logo (if exists)
    logo = os.path.join(IMG, "logo-khtn.png")
    if os.path.exists(logo):
        add_image_fit(slide, logo, Inches(0.6), Inches(5.3), Inches(1.6), Inches(1.6))

    # Right panel
    add_text(slide, Inches(5.0), Inches(1.6), Inches(7.8), Inches(0.35),
             "Đề tài", size=11, bold=True, color=COLOR_ACCENT, font=FONT_TITLE)
    add_text(slide, Inches(5.0), Inches(2.0), Inches(8.0), Inches(2.4),
             "PHÁT HIỆN RANH GIỚI\nCẢNH QUAY DỰA TRÊN\nMÔ HÌNH HỌC SÂU",
             size=34, bold=True, color=COLOR_INK, font=FONT_TITLE)

    add_line(slide, Inches(5.0), Inches(4.55), Inches(7.0), Inches(4.55),
             color=COLOR_ACCENT2, width=2.0)

    add_text(slide, Inches(5.0), Inches(4.75), Inches(7.8), Inches(0.35),
             "Sinh viên thực hiện", size=11, bold=True,
             color=COLOR_SUB, font=FONT_TITLE)
    add_text(slide, Inches(5.0), Inches(5.1), Inches(7.8), Inches(0.4),
             "Tặng Sềnh Mành — 22120202",
             size=16, bold=True, color=COLOR_INK)
    add_text(slide, Inches(5.0), Inches(5.45), Inches(7.8), Inches(0.4),
             "Đỗ Tiến Mạnh — 22120203",
             size=16, bold=True, color=COLOR_INK)

    add_text(slide, Inches(5.0), Inches(6.0), Inches(7.8), Inches(0.35),
             "Giảng viên hướng dẫn", size=11, bold=True,
             color=COLOR_SUB, font=FONT_TITLE)
    add_text(slide, Inches(5.0), Inches(6.35), Inches(7.8), Inches(0.4),
             "TS. Trần Thái Sơn", size=16, bold=True, color=COLOR_INK)

    add_text(slide, Inches(5.0), Inches(6.9), Inches(7.8), Inches(0.3),
             "TP. Hồ Chí Minh  •  Tháng 04/2026",
             size=11, color=COLOR_SUB)


def slide_toc(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Nội dung trình bày", "Mục lục", 2)

    items = [
        ("1", "Giới thiệu", "Bối cảnh, bài toán, đóng góp"),
        ("2", "Công trình liên quan", "Từ cổ điển đến NAS"),
        ("3", "Phương pháp đề xuất", "AutoShotV2"),
        ("4", "Thực nghiệm", "Dữ liệu, kết quả, ablation"),
        ("5", "Kết luận", "Đóng góp, hạn chế, hướng phát triển"),
    ]

    top0 = Inches(2.0)
    row_h = Inches(0.85)
    for i, (num, name, desc) in enumerate(items):
        y = top0 + row_h * i
        # number badge
        badge = add_round_rect(slide, Inches(0.8), y, Inches(0.7), Inches(0.7),
                               fill=COLOR_ACCENT)
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.name = FONT_TITLE
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = COLOR_WHITE

        add_text(slide, Inches(1.8), y + Inches(0.05), Inches(11), Inches(0.45),
                 name, size=22, bold=True, color=COLOR_INK, font=FONT_TITLE)
        add_text(slide, Inches(1.8), y + Inches(0.48), Inches(11), Inches(0.35),
                 desc, size=13, color=COLOR_SUB)


def slide_summary(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Tóm tắt một trang", "Executive summary", 3)

    # 3 KPI cards
    labels = [
        ("853", "video trong tập SHOT", "11.606 nhãn • 960.794 khung hình"),
        (
            f1_text(RESULTS["summary"]["shot_f1"]),
            "F1 của AutoShotV2 trên SHOT",
            "đạt trạng thái tốt nhất",
        ),
        (
            f"+{RESULTS['summary']['shot_delta_vs_transnet_pp']:.1f}%",
            "cải thiện so với TransNetV2",
            f"+{RESULTS['summary']['shot_delta_vs_autoshot_pp']:.1f}% so với AutoShot gốc",
        ),
    ]
    card_w = Inches(3.9)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    top = Inches(2.0)
    for i, (kpi, head, foot) in enumerate(labels):
        x = start_x + (card_w + gap) * i
        card = add_round_rect(slide, x, top, card_w, Inches(2.6),
                              fill=COLOR_PANEL)
        add_text(slide, x, top + Inches(0.4), card_w, Inches(1.0),
                 kpi, size=52, bold=True, color=COLOR_ACCENT,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)
        add_text(slide, x, top + Inches(1.55), card_w, Inches(0.45),
                 head, size=14, bold=True, color=COLOR_INK,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x, top + Inches(2.0), card_w, Inches(0.4),
                 foot, size=11, color=COLOR_SUB, align=PP_ALIGN.CENTER)

    # Tagline
    add_rect(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.4),
             fill=COLOR_INK)
    add_text(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(0.5),
             "Đóng góp cốt lõi", size=12, bold=True,
             color=COLOR_ACCENT2, font=FONT_TITLE)
    add_text(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.9),
             "Kết hợp Focal Loss + Many-hot labels + Temperature scaling "
             "+ Gaussian smoothing trên backbone NAS, đóng băng hoàn toàn "
             "để huấn luyện lại classification head cho video ngắn.",
             size=16, color=COLOR_WHITE)


def slide_motivation(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Bối cảnh & động lực", "Chương 1", 4)

    img = os.path.join(IMG, "tang_truong_short.png")
    if os.path.exists(img):
        add_image_fit(slide, img, Inches(0.6), Inches(1.8),
                      Inches(7.0), Inches(5.0))

    add_bullets(slide, Inches(8.0), Inches(2.0), Inches(4.8), Inches(5.0), [
        ("Short-form bùng nổ: ", "TikTok, Reels, Shorts chiếm lĩnh nội dung số."),
        ("Nhu cầu phân tích tự động: ", "chỉ mục, đề xuất, kiểm duyệt, quảng cáo."),
        ("Shot Boundary Detection (SBD): ", "bước nền cho hầu hết pipeline phân tích video."),
        ("Phương pháp cũ lệch miền: ", "hiệu năng giảm mạnh trên video ngắn."),
    ], size=16)


def slide_problem(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Bài toán Shot Boundary Detection", "Chương 1", 5)

    img = os.path.join(IMG, "video_scene.png")
    if os.path.exists(img):
        add_image_fit(slide, img, Inches(0.6), Inches(1.8),
                      Inches(6.5), Inches(5.0))

    # Definition box
    box = add_round_rect(slide, Inches(7.5), Inches(2.0),
                         Inches(5.3), Inches(2.0), fill=COLOR_PANEL)
    add_text(slide, Inches(7.75), Inches(2.15), Inches(5.0), Inches(0.35),
             "Định nghĩa hình thức", size=11, bold=True,
             color=COLOR_ACCENT, font=FONT_TITLE)
    add_text(slide, Inches(7.75), Inches(2.5), Inches(5.0), Inches(1.4),
             "Input:  V = {f₁, f₂, …, fₙ}\n"
             "Output: B = {b₁, b₂, …, bₘ}\n"
             "              + loại chuyển cảnh",
             size=16, color=COLOR_INK)

    add_text(slide, Inches(7.5), Inches(4.3), Inches(5.3), Inches(0.4),
             "Hai loại chuyển cảnh", size=14, bold=True,
             color=COLOR_INK, font=FONT_TITLE)
    add_bullets(slide, Inches(7.5), Inches(4.75), Inches(5.3), Inches(2.0), [
        ("Cut: ", "ranh giới đột ngột giữa 2 khung hình liên tiếp."),
        ("Gradual: ", "fade, dissolve, wipe — thay đổi trên nhiều khung."),
    ], size=14)


def slide_transitions(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Minh họa các loại chuyển cảnh", "Chương 1", 6)

    transitions = [
        ("cut_transition.png", "Cut",     "Đổi đột ngột giữa 2 khung"),
        ("dissolve.png",       "Dissolve","Hòa trộn tuyến tính"),
        ("fade_in.png",        "Fade",    "Xuất hiện / biến mất dần"),
        ("wipe.png",           "Wipe",    "Quét ranh giới hình học"),
    ]
    cell_w = Inches(2.95)
    cell_h = Inches(2.35)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    top = Inches(2.0)
    for i, (fname, name, desc) in enumerate(transitions):
        x = start_x + (cell_w + gap) * i
        add_round_rect(slide, x, top, cell_w, cell_h, fill=COLOR_PANEL)
        p = os.path.join(IMG, fname)
        if os.path.exists(p):
            add_image_fit(slide, p, x + Inches(0.15), top + Inches(0.15),
                          cell_w - Inches(0.3), cell_h - Inches(0.3))
        # caption area
        add_rect(slide, x, top + cell_h + Inches(0.1), cell_w, Inches(0.9),
                 fill=COLOR_PANEL)
        add_text(slide, x + Inches(0.15), top + cell_h + Inches(0.15),
                 cell_w - Inches(0.3), Inches(0.4),
                 name, size=16, bold=True, color=COLOR_ACCENT,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)
        add_text(slide, x + Inches(0.15), top + cell_h + Inches(0.55),
                 cell_w - Inches(0.3), Inches(0.4),
                 desc, size=11, color=COLOR_SUB, align=PP_ALIGN.CENTER)


def slide_challenges(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Thách thức riêng của video ngắn", "Chương 1", 7)

    # Left: two images stacked
    p1 = os.path.join(IMG, "tiktok.jpg")
    p2 = os.path.join(IMG, "flash.jpg")
    if os.path.exists(p1):
        add_image_fit(slide, p1, Inches(0.6), Inches(1.85),
                      Inches(5.0), Inches(2.55))
    if os.path.exists(p2):
        add_image_fit(slide, p2, Inches(0.6), Inches(4.5),
                      Inches(5.0), Inches(2.35))

    items = [
        ("Shot dày đặc: ", "trung bình 2,59 s / shot (vs 6–15 s truyền thống)."),
        ("Hiệu ứng phức tạp: ", "gradual chồng lớp, hiệu ứng editor đa dạng."),
        ("Khung dọc có viền tĩnh: ", "cấu trúc ba phần gây nhầm lẫn cho SBD."),
        ("Flash & camera motion: ", "tạo false positive khó phân biệt."),
        ("Nội dung ảo / game: ", "biến đổi thị giác ngay trong một shot."),
    ]
    add_bullets(slide, Inches(5.9), Inches(1.85), Inches(7.0), Inches(5.0),
                items, size=16, line_spacing=1.4)


def slide_rq_contrib(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Câu hỏi nghiên cứu & đóng góp", "Chương 1", 8)

    # Two columns
    col_w = Inches(6.0)
    gap = Inches(0.2)

    # Left: research questions
    add_rect(slide, Inches(0.6), Inches(1.85), col_w, Inches(0.4),
             fill=COLOR_ACCENT)
    add_text(slide, Inches(0.8), Inches(1.9), col_w, Inches(0.35),
             "CÂU HỎI NGHIÊN CỨU", size=12, bold=True,
             color=COLOR_WHITE, font=FONT_TITLE)
    add_round_rect(slide, Inches(0.6), Inches(2.3), col_w, Inches(4.6),
                   fill=COLOR_PANEL)
    rqs = [
        ("RQ1", "NAS có giúp kiến trúc vượt trội so với backbone cố định trên video ngắn?"),
        ("RQ2", "Các kỹ thuật head + hậu xử lý (Focal, many-hot, temp scaling, Gaussian) có cải thiện F1 mà không đổi backbone?"),
        ("RQ3", "Cải tiến trên SHOT có tổng quát hóa sang BBC / ClipShots?"),
    ]
    y = Inches(2.5)
    for tag, q in rqs:
        add_round_rect(slide, Inches(0.85), y, Inches(0.9), Inches(0.55),
                       fill=COLOR_ACCENT)
        add_text(slide, Inches(0.85), y + Inches(0.08), Inches(0.9), Inches(0.4),
                 tag, size=13, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)
        add_text(slide, Inches(1.9), y, col_w - Inches(1.4), Inches(1.3),
                 q, size=14, color=COLOR_INK)
        y += Inches(1.4)

    # Right: contributions
    x2 = Inches(0.6) + col_w + gap
    add_rect(slide, x2, Inches(1.85), col_w, Inches(0.4),
             fill=COLOR_ACCENT2)
    add_text(slide, x2 + Inches(0.2), Inches(1.9), col_w, Inches(0.35),
             "ĐÓNG GÓP CHÍNH", size=12, bold=True,
             color=COLOR_WHITE, font=FONT_TITLE)
    add_round_rect(slide, x2, Inches(2.3), col_w, Inches(4.6),
                   fill=COLOR_PANEL)
    contribs = [
        ("01", "Tổng quan có hệ thống", "SBD cổ điển → deep learning → NAS."),
        ("02", "Pipeline AutoShotV2", "Focal Loss, many-hot, temp scaling, Gaussian trên backbone NAS đóng băng."),
        ("03", "Đánh giá đa miền", "SHOT, BBC, ClipShots — phân tích tổng quát hóa."),
    ]
    y = Inches(2.5)
    for n, name, desc in contribs:
        add_text(slide, x2 + Inches(0.25), y, Inches(1.0), Inches(0.5),
                 n, size=26, bold=True, color=COLOR_ACCENT2, font=FONT_TITLE)
        add_text(slide, x2 + Inches(1.1), y, col_w - Inches(1.3), Inches(0.4),
                 name, size=15, bold=True, color=COLOR_INK, font=FONT_TITLE)
        add_text(slide, x2 + Inches(1.1), y + Inches(0.45),
                 col_w - Inches(1.3), Inches(0.9),
                 desc, size=12, color=COLOR_SUB)
        y += Inches(1.4)


def slide_sbd_history(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Tiến hóa các phương pháp SBD", "Chương 2", 9)

    phases = [
        ("Cổ điển",
         "Pre-2016",
         "Histogram • pixel diff • edge",
         "Nhẹ, trực quan nhưng kém tổng quát hóa."),
        ("Deep Learning",
         "2017 – 2020",
         "C3D (DeepSBD) • TransNetV2",
         "Học đặc trưng không – thời gian; mạnh trên phim dài."),
        ("NAS",
         "2023 →",
         "AutoShot • AutoShotV2",
         "Tự tìm kiến trúc tối ưu cho short-form video."),
    ]
    cell_w = Inches(4.0)
    cell_h = Inches(4.4)
    gap = Inches(0.15)
    start_x = Inches(0.6)
    top = Inches(2.0)
    for i, (name, years, methods, note) in enumerate(phases):
        x = start_x + (cell_w + gap) * i
        add_round_rect(slide, x, top, cell_w, cell_h, fill=COLOR_PANEL)
        # colored header strip
        color = [COLOR_SUB, COLOR_ACCENT, COLOR_ACCENT2][i]
        add_rect(slide, x, top, cell_w, Inches(0.5), fill=color)
        add_text(slide, x + Inches(0.25), top + Inches(0.1),
                 cell_w - Inches(0.5), Inches(0.3),
                 years, size=11, bold=True, color=COLOR_WHITE, font=FONT_TITLE)
        add_text(slide, x + Inches(0.25), top + Inches(0.7),
                 cell_w - Inches(0.5), Inches(0.5),
                 name, size=22, bold=True, color=COLOR_INK, font=FONT_TITLE)
        add_line(slide, x + Inches(0.25), top + Inches(1.3),
                 x + Inches(1.5), top + Inches(1.3),
                 color=color, width=2.5)
        add_text(slide, x + Inches(0.25), top + Inches(1.55),
                 cell_w - Inches(0.5), Inches(0.9),
                 methods, size=14, bold=True, color=COLOR_INK)
        add_text(slide, x + Inches(0.25), top + Inches(2.55),
                 cell_w - Inches(0.5), Inches(1.8),
                 note, size=13, color=COLOR_SUB)


def slide_baselines(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "TransNetV2 & AutoShot — baseline", "Chương 2", 10)

    # Comparison table
    headers = ["Mô hình", "F1 SHOT", "F1 BBC", "F1 ClipShots", "Ghi chú"]
    rows = [
        (
            "TransNetV2",
            f1_text(comparison_f1("transnetv2_reported", "shot")),
            f1_text(comparison_f1("transnetv2_reported", "bbc")),
            f1_text(comparison_f1("transnetv2_reported", "clipshots")),
            "Tối ưu cho phim dài",
        ),
        (
            "AutoShot",
            f1_text(comparison_f1("autoshot_reported", "shot")),
            f1_text(comparison_f1("autoshot_reported", "bbc")),
            f1_text(comparison_f1("autoshot_reported", "clipshots")),
            "NAS cho short-form",
        ),
    ]
    table_top = Inches(2.0)
    col_widths = [Inches(2.2), Inches(1.5), Inches(1.5), Inches(1.9), Inches(3.0)]
    start_x = Inches(0.6)
    row_h = Inches(0.55)

    # header row
    x = start_x
    add_rect(slide, x, table_top, sum(col_widths, Emu(0)), row_h,
             fill=COLOR_ACCENT)
    for w, h in zip(col_widths, headers):
        add_text(slide, x + Inches(0.15), table_top + Inches(0.12),
                 w - Inches(0.3), row_h,
                 h, size=13, bold=True, color=COLOR_WHITE, font=FONT_TITLE)
        x += w

    # data rows
    for i, row in enumerate(rows):
        x = start_x
        y = table_top + row_h * (i + 1)
        bg = COLOR_PANEL if i % 2 == 0 else COLOR_WHITE
        add_rect(slide, x, y, sum(col_widths, Emu(0)), row_h, fill=bg)
        for j, (w, val) in enumerate(zip(col_widths, row)):
            bold = (j == 0)
            add_text(slide, x + Inches(0.15), y + Inches(0.12),
                     w - Inches(0.3), row_h,
                     val, size=14, bold=bold, color=COLOR_INK)
            x += w

    # Key observations
    add_text(slide, Inches(0.6), Inches(4.4), Inches(12.1), Inches(0.4),
             "Nhận xét", size=14, bold=True, color=COLOR_ACCENT, font=FONT_TITLE)
    add_bullets(slide, Inches(0.6), Inches(4.85), Inches(12.1), Inches(2.0), [
        ("TransNetV2 tối ưu cho phim dài: ",
         "F1 giảm khi chuyển sang video ngắn."),
        ("AutoShot (NAS) vượt baseline ", "+4,2% trên SHOT — cho thấy NAS phù hợp với miền short-form."),
        ("Vẫn còn dư địa: ",
         "head classification và hậu xử lý chưa được khai thác hệ thống."),
    ], size=15)


def slide_gaps(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Khoảng trống nghiên cứu", "Chương 2", 11)

    gaps = [
        ("Domain shift",
         "Bộ dữ liệu truyền thống (BBC, RAI) có đặc trưng khác hẳn short-form video hiện đại."),
        ("Tổng quát hóa đa miền",
         "F1 cao trên SHOT không đảm bảo F1 cao trên BBC hay ClipShots."),
        ("Head & post-processing",
         "Ít nghiên cứu hệ thống về classification head + hiệu chỉnh xác suất + làm mượt thời gian."),
    ]
    top = Inches(2.2)
    row_h = Inches(1.45)
    for i, (name, desc) in enumerate(gaps):
        y = top + row_h * i
        add_round_rect(slide, Inches(0.6), y, Inches(1.3), Inches(1.3),
                       fill=COLOR_ACCENT)
        add_text(slide, Inches(0.6), y + Inches(0.35), Inches(1.3), Inches(0.6),
                 f"0{i+1}", size=32, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)
        add_text(slide, Inches(2.1), y + Inches(0.1), Inches(10.6), Inches(0.45),
                 name, size=20, bold=True, color=COLOR_INK, font=FONT_TITLE)
        add_text(slide, Inches(2.1), y + Inches(0.55), Inches(10.6), Inches(0.9),
                 desc, size=14, color=COLOR_SUB)

    add_text(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
             "→ AutoShotV2 tập trung lấp khoảng trống thứ ba, đồng thời đo kiểm tổng quát hóa.",
             size=14, bold=True, color=COLOR_ACCENT2, font=FONT_TITLE)


def slide_arch_overview(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Kiến trúc tổng quan AutoShot / AutoShotV2", "Chương 3", 12)

    img = os.path.join(IMG, "Flowchart.png")
    if os.path.exists(img):
        add_image_fit(slide, img, Inches(0.6), Inches(1.8),
                      Inches(8.2), Inches(5.2))

    add_text(slide, Inches(9.0), Inches(2.0), Inches(3.9), Inches(0.4),
             "Pipeline 4 giai đoạn", size=13, bold=True,
             color=COLOR_ACCENT, font=FONT_TITLE)
    add_bullets(slide, Inches(9.0), Inches(2.5), Inches(3.9), Inches(4.8), [
        ("Input: ", "chuỗi khung hình video."),
        ("Backbone: ", "3D CNN + Transformer (tìm bằng NAS)."),
        ("Head: ", "dự đoán xác suất boundary."),
        ("Post-proc: ", "smoothing + ngưỡng hóa → boundary."),
    ], size=14, line_spacing=1.35)


def slide_nas(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Neural Architecture Search", "Chương 3", 13)

    # Left: key numbers
    add_round_rect(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(5.0),
                   fill=COLOR_PANEL)
    add_text(slide, Inches(0.8), Inches(2.15), Inches(5.4), Inches(0.4),
             "Không gian tìm kiếm", size=13, bold=True,
             color=COLOR_ACCENT, font=FONT_TITLE)
    add_text(slide, Inches(0.8), Inches(2.55), Inches(5.4), Inches(1.2),
             "83,9 triệu", size=56, bold=True, color=COLOR_ACCENT,
             font=FONT_TITLE)
    add_text(slide, Inches(0.8), Inches(3.85), Inches(5.4), Inches(0.4),
             "kiến trúc ứng viên", size=14, color=COLOR_INK)
    add_text(slide, Inches(0.8), Inches(4.3), Inches(5.4), Inches(0.4),
             "= 4 biến thể DDCNNV2 × 6 block stack",
             size=12, color=COLOR_SUB)
    add_line(slide, Inches(0.8), Inches(4.9), Inches(6.2), Inches(4.9),
             color=COLOR_LINE, width=0.75)
    add_bullets(slide, Inches(0.8), Inches(5.05), Inches(5.4), Inches(2.0), [
        ("Hybrid: ", "3D CNN + Transformer block."),
        ("DDCNNV2: ", "dilated conv + kernel factorization."),
    ], size=13, line_spacing=1.3)

    # Right: three phases
    x = Inches(6.8)
    add_text(slide, x, Inches(2.0), Inches(6.0), Inches(0.4),
             "Ba thành phần của NAS", size=13, bold=True,
             color=COLOR_ACCENT, font=FONT_TITLE)
    phases = [
        ("Search Space",
         "Tập hợp kiến trúc có thể tạo ra (83,9 M)."),
        ("Search Strategy",
         "One-shot SuperNet: chia sẻ trọng số giữa các ứng viên."),
        ("Performance Estimation",
         "Early stopping + gradient-based refinement để đánh giá nhanh."),
    ]
    y = Inches(2.55)
    for name, desc in phases:
        add_round_rect(slide, x, y, Inches(6.0), Inches(1.35),
                       fill=COLOR_PANEL)
        add_rect(slide, x, y, Inches(0.2), Inches(1.35), fill=COLOR_ACCENT)
        add_text(slide, x + Inches(0.4), y + Inches(0.15),
                 Inches(5.4), Inches(0.45),
                 name, size=15, bold=True, color=COLOR_INK, font=FONT_TITLE)
        add_text(slide, x + Inches(0.4), y + Inches(0.6),
                 Inches(5.4), Inches(0.75),
                 desc, size=12, color=COLOR_SUB)
        y += Inches(1.5)


def slide_loss(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "AutoShotV2 — cải tiến hàm mất mát & nhãn", "Chương 3", 14)

    col_w = Inches(6.0)
    # Focal loss box
    add_round_rect(slide, Inches(0.6), Inches(1.9), col_w, Inches(4.9),
                   fill=COLOR_PANEL)
    add_rect(slide, Inches(0.6), Inches(1.9), col_w, Inches(0.55),
             fill=COLOR_ACCENT)
    add_text(slide, Inches(0.85), Inches(2.0), col_w, Inches(0.4),
             "FOCAL LOSS", size=14, bold=True,
             color=COLOR_WHITE, font=FONT_TITLE)
    add_text(slide, Inches(0.85), Inches(2.65), col_w - Inches(0.4), Inches(1.2),
             "FL(p) = −α (1 − p)^γ log(p)",
             size=22, bold=True, color=COLOR_INK, font=FONT_TITLE)
    add_text(slide, Inches(0.85), Inches(3.55), col_w - Inches(0.4), Inches(0.5),
             "γ = 1.0, α = 0.5   (grid search)",
             size=12, color=COLOR_SUB)
    add_bullets(slide, Inches(0.85), Inches(4.2), col_w - Inches(0.4), Inches(2.6), [
        "Giảm trọng số cho khung non-boundary (dễ).",
        "Tập trung học trên mẫu hard / boundary hiếm.",
        "Giải quyết mất cân bằng lớp nặng của SBD.",
    ], size=14, line_spacing=1.3)

    # Many-hot box
    x2 = Inches(0.6) + col_w + Inches(0.3)
    add_round_rect(slide, x2, Inches(1.9), col_w, Inches(4.9),
                   fill=COLOR_PANEL)
    add_rect(slide, x2, Inches(1.9), col_w, Inches(0.55),
             fill=COLOR_ACCENT2)
    add_text(slide, x2 + Inches(0.25), Inches(2.0), col_w, Inches(0.4),
             "MANY-HOT LABELS", size=14, bold=True,
             color=COLOR_WHITE, font=FONT_TITLE)
    add_text(slide, x2 + Inches(0.25), Inches(2.65), col_w - Inches(0.4),
             Inches(1.2),
             "[…, 0, 0.5, 1.0, 0.5, 0, …]",
             size=20, bold=True, color=COLOR_INK, font=FONT_TITLE)
    add_text(slide, x2 + Inches(0.25), Inches(3.55),
             col_w - Inches(0.4), Inches(0.5),
             "thay vì {0, 1} cứng tại ranh giới",
             size=12, color=COLOR_SUB)
    add_bullets(slide, x2 + Inches(0.25), Inches(4.2),
                col_w - Inches(0.4), Inches(2.6), [
        "Nhãn mềm quanh biên giúp mô hình học chuyển tiếp mượt.",
        "Giảm penalty khi dự đoán lệch 1 khung.",
        "Kết hợp song song với one-hot trong head hai đầu ra.",
    ], size=14, line_spacing=1.3)


def slide_calibration(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Hiệu chỉnh xác suất & hậu xử lý", "Chương 3", 15)

    col_w = Inches(6.0)
    # Temperature scaling
    add_round_rect(slide, Inches(0.6), Inches(1.9), col_w, Inches(4.9),
                   fill=COLOR_PANEL)
    add_rect(slide, Inches(0.6), Inches(1.9), col_w, Inches(0.55),
             fill=COLOR_ACCENT)
    add_text(slide, Inches(0.85), Inches(2.0), col_w, Inches(0.4),
             "TEMPERATURE SCALING", size=14, bold=True,
             color=COLOR_WHITE, font=FONT_TITLE)
    add_text(slide, Inches(0.85), Inches(2.7), col_w - Inches(0.4), Inches(1.0),
             "p = σ(z / T)",
             size=28, bold=True, color=COLOR_INK, font=FONT_TITLE)
    add_text(slide, Inches(0.85), Inches(3.7), col_w - Inches(0.4), Inches(0.5),
             "Học T trên validation để confidence khớp tần suất đúng.",
             size=12, color=COLOR_SUB)
    add_bullets(slide, Inches(0.85), Inches(4.3), col_w - Inches(0.4), Inches(2.4), [
        ("Ngưỡng tối ưu dịch: ", "0.296 → 0.02."),
        "Phân bố xác suất nhất quán hơn giữa các video.",
        "Dễ dàng hơn cho bước ngưỡng hóa.",
    ], size=14, line_spacing=1.3)

    # Gaussian smoothing
    x2 = Inches(0.6) + col_w + Inches(0.3)
    add_round_rect(slide, x2, Inches(1.9), col_w, Inches(4.9),
                   fill=COLOR_PANEL)
    add_rect(slide, x2, Inches(1.9), col_w, Inches(0.55),
             fill=COLOR_ACCENT2)
    add_text(slide, x2 + Inches(0.25), Inches(2.0), col_w, Inches(0.4),
             "GAUSSIAN SMOOTHING", size=14, bold=True,
             color=COLOR_WHITE, font=FONT_TITLE)
    add_text(slide, x2 + Inches(0.25), Inches(2.7), col_w - Inches(0.4), Inches(1.0),
             "p̃ = Gaussian(p; σ = 2.0)",
             size=22, bold=True, color=COLOR_INK, font=FONT_TITLE)
    add_text(slide, x2 + Inches(0.25), Inches(3.7),
             col_w - Inches(0.4), Inches(0.5),
             "Tích chập 1-D trên chuỗi xác suất thời gian.",
             size=12, color=COLOR_SUB)
    add_bullets(slide, x2 + Inches(0.25), Inches(4.3),
                col_w - Inches(0.4), Inches(2.4), [
        "Triệt tiêu false positive đơn lẻ.",
        "Giữ vùng đỉnh liên tục quanh boundary thật.",
        ("Đóng góp rõ rệt nhất: ", "+3.02% trên ClipShots."),
    ], size=14, line_spacing=1.3)


def slide_training(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Chiến lược huấn luyện", "Chương 3", 16)

    # Two columns: strategy & hyperparameters
    # Left: strategy
    add_rect(slide, Inches(0.6), Inches(1.9), Inches(6.2), Inches(0.45),
             fill=COLOR_ACCENT)
    add_text(slide, Inches(0.8), Inches(1.95), Inches(6.0), Inches(0.35),
             "Chiến lược", size=13, bold=True, color=COLOR_WHITE,
             font=FONT_TITLE)
    add_round_rect(slide, Inches(0.6), Inches(2.35), Inches(6.2), Inches(4.55),
                   fill=COLOR_PANEL)
    add_bullets(slide, Inches(0.85), Inches(2.55), Inches(5.9), Inches(4.3), [
        ("Backbone đóng băng: ", "giữ nguyên trọng số AutoShot."),
        ("Chỉ huấn luyện ClassificationHead ", "(input 4864-D)."),
        ("Warm-up + cosine annealing ", "learning rate."),
        ("EMA ", "cho loss và trọng số → ổn định."),
        ("Early stopping ", "theo F1 trên validation."),
        ("Focal Loss + Many-hot labels ", "kết hợp trong head 2 đầu ra."),
    ], size=14, line_spacing=1.35)

    # Right: hyperparameters table
    x2 = Inches(7.0)
    add_rect(slide, x2, Inches(1.9), Inches(5.7), Inches(0.45),
             fill=COLOR_ACCENT2)
    add_text(slide, x2 + Inches(0.2), Inches(1.95), Inches(5.3), Inches(0.35),
             "Siêu tham số", size=13, bold=True, color=COLOR_WHITE,
             font=FONT_TITLE)
    params = [
        ("Focal Loss γ",  "1.0"),
        ("Focal Loss α",  "0.5"),
        ("Gaussian σ",    "2.0"),
        ("Threshold (V2)", "0.02"),
        ("LR schedule",   "Warm-up + Cosine"),
        ("EMA",           "Bật"),
        ("Head size",     "4864 → 1024 → 1"),
    ]
    row_h = Inches(0.55)
    for i, (k, v) in enumerate(params):
        y = Inches(2.45) + row_h * i
        bg = COLOR_PANEL if i % 2 == 0 else COLOR_WHITE
        add_rect(slide, x2, y, Inches(5.7), row_h, fill=bg)
        add_text(slide, x2 + Inches(0.25), y + Inches(0.12),
                 Inches(3.0), Inches(0.4),
                 k, size=13, color=COLOR_INK)
        add_text(slide, x2 + Inches(3.3), y + Inches(0.12),
                 Inches(2.3), Inches(0.4),
                 v, size=13, bold=True, color=COLOR_ACCENT, font=FONT_TITLE)


def slide_diff(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "AutoShot vs AutoShotV2", "Chương 3", 17)

    headers = ["Thành phần", "AutoShot", "AutoShotV2"]
    rows = [
        ("Backbone NAS",          "Huấn luyện full",            "Đóng băng hoàn toàn"),
        ("Classification head",   "Lớp gốc, fc1_0",             "Head mới (4864 → 1024 → 1)"),
        ("Loss function",         "Binary cross-entropy",       "Focal Loss (γ=1.0, α=0.5)"),
        ("Label format",          "One-hot",                    "One-hot + Many-hot"),
        ("Calibration",           "—",                          "Temperature scaling"),
        ("Post-processing",       "—",                          "Gaussian smoothing (σ=2.0)"),
        ("F1 trên SHOT",          "0.841",                      "0.855  (+1.4%)"),
    ]
    col_widths = [Inches(3.8), Inches(4.2), Inches(4.1)]
    start_x = Inches(0.6)
    row_h = Inches(0.55)
    table_top = Inches(2.0)

    # header
    x = start_x
    add_rect(slide, x, table_top, sum(col_widths, Emu(0)), row_h,
             fill=COLOR_INK)
    for w, h in zip(col_widths, headers):
        add_text(slide, x + Inches(0.15), table_top + Inches(0.12),
                 w - Inches(0.3), row_h,
                 h, size=13, bold=True, color=COLOR_WHITE, font=FONT_TITLE)
        x += w

    # rows
    for i, row in enumerate(rows):
        x = start_x
        y = table_top + row_h * (i + 1)
        bg = COLOR_PANEL if i % 2 == 0 else COLOR_WHITE
        highlight = (i == len(rows) - 1)
        if highlight:
            bg = RGBColor(0xE7, 0xF0, 0xFF)
        add_rect(slide, x, y, sum(col_widths, Emu(0)), row_h, fill=bg)
        for j, (w, val) in enumerate(zip(col_widths, row)):
            color = COLOR_INK if j == 0 else COLOR_INK
            bold = (j == 0) or highlight
            if highlight and j == 2:
                color = COLOR_ACCENT
            add_text(slide, x + Inches(0.15), y + Inches(0.12),
                     w - Inches(0.3), row_h,
                     val, size=13, bold=bold, color=color)
            x += w


def slide_dataset_shot(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Bộ dữ liệu SHOT", "Chương 4", 18)

    img = os.path.join(IMG, "Anh41.png")
    if os.path.exists(img):
        add_image_fit(slide, img, Inches(0.6), Inches(1.9),
                      Inches(7.0), Inches(5.0))

    # Right: statistics cards
    stats = [
        ("853",     "video"),
        ("11.606",  "nhãn boundary"),
        ("960.794", "khung hình"),
        ("39,5 s",  "độ dài video TB"),
        ("2,59 s",  "độ dài shot TB"),
    ]
    top = Inches(2.0)
    row_h = Inches(0.95)
    for i, (num, label) in enumerate(stats):
        y = top + row_h * i
        add_round_rect(slide, Inches(7.9), y, Inches(5.0), Inches(0.85),
                       fill=COLOR_PANEL)
        add_text(slide, Inches(8.1), y + Inches(0.15), Inches(2.0), Inches(0.55),
                 num, size=22, bold=True, color=COLOR_ACCENT, font=FONT_TITLE)
        add_text(slide, Inches(10.2), y + Inches(0.22),
                 Inches(2.6), Inches(0.5),
                 label, size=14, color=COLOR_INK)


def slide_dataset_challenges(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Đặc trưng khó của SHOT", "Chương 4", 19)

    img = os.path.join(IMG, "Anh42.png")
    if os.path.exists(img):
        add_image_fit(slide, img, Inches(0.6), Inches(1.9),
                      Inches(8.0), Inches(5.0))

    add_bullets(slide, Inches(8.9), Inches(2.0), Inches(4.0), Inches(5.0), [
        ("Gradual phức tạp: ", "dissolve / wipe chồng nhiều lớp."),
        ("Cấu trúc ba phần: ", "viền tĩnh trên–dưới gây nhiễu."),
        ("Game / ảo: ", "biến đổi nội bộ shot mạnh."),
        ("Camera motion: ", "pan/zoom nhanh giống boundary."),
        ("Flash & lighting: ", "thay đổi đột ngột — false positive."),
    ], size=14, line_spacing=1.4)


def slide_dataset_others(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Các bộ dữ liệu đối chứng", "Chương 4", 20)

    headers = ["Dataset", "Số video", "Độ dài video TB", "Độ dài shot TB", "Đặc trưng"]
    rows = [
        ("SHOT",           "853",  "39,5 s",  "2,59 s",  "Short-form đa nền tảng"),
        ("BBC Planet Earth","11",  "≈ 49 min","6,57 s",  "Phim tài liệu thiên nhiên"),
        ("RAI",            "10",   "≈ 10 min","5,65 s",  "TV giáo dục Ý"),
        ("ClipShots",      "~500", "237 s",   "15,34 s", "YouTube/Weibo đa dạng"),
    ]
    col_widths = [Inches(1.8), Inches(1.6), Inches(2.4), Inches(2.2), Inches(4.1)]
    start_x = Inches(0.6)
    row_h = Inches(0.6)
    table_top = Inches(2.0)

    # header
    x = start_x
    add_rect(slide, x, table_top, sum(col_widths, Emu(0)), row_h,
             fill=COLOR_ACCENT)
    for w, h in zip(col_widths, headers):
        add_text(slide, x + Inches(0.15), table_top + Inches(0.15),
                 w - Inches(0.3), row_h,
                 h, size=13, bold=True, color=COLOR_WHITE, font=FONT_TITLE)
        x += w

    for i, row in enumerate(rows):
        x = start_x
        y = table_top + row_h * (i + 1)
        bg = COLOR_PANEL if i % 2 == 0 else COLOR_WHITE
        if row[0] == "SHOT":
            bg = RGBColor(0xE7, 0xF0, 0xFF)
        add_rect(slide, x, y, sum(col_widths, Emu(0)), row_h, fill=bg)
        for j, (w, val) in enumerate(zip(col_widths, row)):
            bold = (j == 0) or (row[0] == "SHOT")
            color = COLOR_ACCENT if (row[0] == "SHOT" and j == 0) else COLOR_INK
            add_text(slide, x + Inches(0.15), y + Inches(0.17),
                     w - Inches(0.3), row_h,
                     val, size=13, bold=bold, color=color)
            x += w

    add_text(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.4),
             "Nhận xét", size=13, bold=True, color=COLOR_ACCENT, font=FONT_TITLE)
    add_bullets(slide, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.3), [
        "SHOT có shot ngắn gấp 2–6 lần so với các bộ truyền thống → miền hoàn toàn khác.",
        "Đánh giá trên cả 4 bộ giúp đo kiểm khả năng tổng quát hóa của mô hình.",
    ], size=13, line_spacing=1.3)


def slide_metrics(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Độ đo đánh giá", "Chương 4", 21)

    metrics = [
        ("Precision", "TP / (TP + FP)",
         "Tỉ lệ dự đoán đúng trong số boundary dự đoán."),
        ("Recall",    "TP / (TP + FN)",
         "Tỉ lệ boundary thật được phát hiện."),
        ("F1-Score",  "2 · P · R / (P + R)",
         "Trung bình điều hòa — độ đo chính của SBD."),
    ]
    cell_w = Inches(4.0)
    cell_h = Inches(3.4)
    gap = Inches(0.15)
    start_x = Inches(0.6)
    top = Inches(2.0)
    for i, (name, formula, desc) in enumerate(metrics):
        x = start_x + (cell_w + gap) * i
        add_round_rect(slide, x, top, cell_w, cell_h, fill=COLOR_PANEL)
        add_rect(slide, x, top, cell_w, Inches(0.6),
                 fill=COLOR_ACCENT if i < 2 else COLOR_ACCENT2)
        add_text(slide, x + Inches(0.25), top + Inches(0.15),
                 cell_w - Inches(0.5), Inches(0.4),
                 name, size=16, bold=True, color=COLOR_WHITE, font=FONT_TITLE)
        add_text(slide, x + Inches(0.25), top + Inches(0.85),
                 cell_w - Inches(0.5), Inches(0.8),
                 formula, size=20, bold=True, color=COLOR_INK, font=FONT_TITLE)
        add_line(slide, x + Inches(0.25), top + Inches(1.85),
                 x + cell_w - Inches(0.25), top + Inches(1.85),
                 color=COLOR_LINE, width=0.75)
        add_text(slide, x + Inches(0.25), top + Inches(2.0),
                 cell_w - Inches(0.5), Inches(1.3),
                 desc, size=13, color=COLOR_SUB)

    # Extra metrics strip
    add_round_rect(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.3),
                   fill=COLOR_INK)
    add_text(slide, Inches(0.85), Inches(5.85), Inches(11.7), Inches(0.4),
             "Độ đo bổ sung", size=12, bold=True,
             color=COLOR_ACCENT2, font=FONT_TITLE)
    add_text(slide, Inches(0.85), Inches(6.25), Inches(11.7), Inches(0.75),
             "Precision @ Fixed Recall  •  F1 riêng cho Cut vs Gradual  •  Đánh giá chéo SHOT / BBC / ClipShots",
             size=14, color=COLOR_WHITE)


def slide_results_main(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Kết quả chính trên SHOT", "Chương 4", 22)

    # Left: table
    headers = ["Mô hình", "F1"]
    rows = [
        ("TransNetV2", f1_text(comparison_f1("transnetv2_reported", "shot")), False),
        ("AutoShot", f1_text(comparison_f1("autoshot_reported", "shot")), False),
        (
            "AutoShotV1 (+Gaussian)",
            f1_text(comparison_f1("autoshot_v1_gaussian", "shot")),
            False,
        ),
        ("AutoShotV2 (3 kỹ thuật)", f1_text(RESULTS["summary"]["shot_f1"]), True),
    ]
    col_widths = [Inches(4.7), Inches(1.8)]
    start_x = Inches(0.6)
    row_h = Inches(0.6)
    table_top = Inches(2.0)

    x = start_x
    add_rect(slide, x, table_top, sum(col_widths, Emu(0)), row_h,
             fill=COLOR_ACCENT)
    for w, h in zip(col_widths, headers):
        add_text(slide, x + Inches(0.15), table_top + Inches(0.15),
                 w - Inches(0.3), row_h,
                 h, size=14, bold=True, color=COLOR_WHITE, font=FONT_TITLE)
        x += w

    for i, (name, f1, hi) in enumerate(rows):
        x = start_x
        y = table_top + row_h * (i + 1)
        bg = RGBColor(0xE7, 0xF0, 0xFF) if hi else (COLOR_PANEL if i % 2 == 0 else COLOR_WHITE)
        add_rect(slide, x, y, sum(col_widths, Emu(0)), row_h, fill=bg)
        add_text(slide, x + Inches(0.15), y + Inches(0.17),
                 col_widths[0] - Inches(0.3), row_h,
                 name, size=14, bold=hi, color=COLOR_INK)
        add_text(slide, x + col_widths[0] + Inches(0.15), y + Inches(0.17),
                 col_widths[1] - Inches(0.3), row_h,
                 f1, size=14, bold=True,
                 color=COLOR_ACCENT if hi else COLOR_INK, font=FONT_TITLE)

    # Right: horizontal bar chart (drawn with shapes)
    chart_x = Inches(7.3)
    chart_y = Inches(2.0)
    chart_w = Inches(5.4)
    chart_h = Inches(3.5)
    add_round_rect(slide, chart_x, chart_y, chart_w, chart_h,
                   fill=COLOR_PANEL)
    add_text(slide, chart_x + Inches(0.25), chart_y + Inches(0.15),
             chart_w - Inches(0.5), Inches(0.3),
             "F1-score trên SHOT", size=11, bold=True,
             color=COLOR_SUB, font=FONT_TITLE)

    # Bars: map F1 from [0.70, 0.90] to [0, chart_w - 1.6"]
    def bar_x_end(value):
        frac = (value - 0.70) / (0.90 - 0.70)
        return chart_x + Inches(1.4) + Inches(frac * 3.6)

    bars = [
        ("TransNetV2", comparison_f1("transnetv2_reported", "shot"), COLOR_SUB),
        ("AutoShot", comparison_f1("autoshot_reported", "shot"), COLOR_ACCENT),
        ("AutoShotV2", RESULTS["summary"]["shot_f1"], COLOR_ACCENT2),
    ]
    bar_h = Inches(0.5)
    bar_gap = Inches(0.25)
    by = chart_y + Inches(0.7)
    for name, val, color in bars:
        # label
        add_text(slide, chart_x + Inches(0.2), by + Inches(0.12),
                 Inches(1.2), bar_h,
                 name, size=11, bold=True, color=COLOR_INK)
        # bar
        end_x = bar_x_end(val)
        bar = add_rect(slide, chart_x + Inches(1.4), by,
                       end_x - chart_x - Inches(1.4), bar_h, fill=color)
        # value
        add_text(slide, end_x + Inches(0.1), by + Inches(0.12),
                 Inches(0.8), bar_h,
                 f"{val:.3f}", size=12, bold=True,
                 color=COLOR_INK, font=FONT_TITLE)
        by += bar_h + bar_gap

    # axis
    axis_y = chart_y + chart_h - Inches(0.55)
    add_line(slide, chart_x + Inches(1.4), axis_y,
             chart_x + Inches(5.0), axis_y,
             color=COLOR_SUB, width=0.75)
    for tick, lbl in [(0.70, "0.70"), (0.75, "0.75"),
                      (0.80, "0.80"), (0.85, "0.85"), (0.90, "0.90")]:
        tx = bar_x_end(tick)
        add_line(slide, tx, axis_y, tx, axis_y + Inches(0.08),
                 color=COLOR_SUB, width=0.5)
        add_text(slide, tx - Inches(0.25), axis_y + Inches(0.1),
                 Inches(0.5), Inches(0.3),
                 lbl, size=9, color=COLOR_SUB, align=PP_ALIGN.CENTER)

    # Takeaway strip
    add_rect(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0),
             fill=COLOR_INK)
    add_text(slide, Inches(0.85), Inches(6.15), Inches(11.7), Inches(0.35),
             "Kết luận", size=11, bold=True,
             color=COLOR_ACCENT2, font=FONT_TITLE)
    add_text(slide, Inches(0.85), Inches(6.5), Inches(11.7), Inches(0.45),
             f"AutoShotV2 đạt F1 = {f1_text(RESULTS['summary']['shot_f1'])} — cao nhất trên SHOT, "
             f"+{RESULTS['summary']['shot_delta_vs_autoshot_pp']:.1f}% so với AutoShot gốc và "
             f"+{RESULTS['summary']['shot_delta_vs_transnet_pp']:.1f}% so với TransNetV2.",
             size=14, bold=True, color=COLOR_WHITE)


def slide_results_cross(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Kết quả trên BBC & ClipShots", "Chương 4", 23)

    headers = ["Mô hình", "SHOT", "BBC", "ClipShots"]
    rows = [
        tuple(
            ["TransNetV2"]
            + [f1_text(comparison_f1("transnetv2_reported", dataset)) for dataset in ("shot", "bbc", "clipshots")]
            + [False]
        ),
        tuple(
            ["AutoShot (báo cáo)"]
            + [f1_text(comparison_f1("autoshot_reported", dataset)) for dataset in ("shot", "bbc", "clipshots")]
            + [False]
        ),
        tuple(
            ["AutoShot tự chạy"]
            + [f1_text(comparison_f1("autoshot_reproduced_legacy", dataset)) for dataset in ("shot", "bbc", "clipshots")]
            + [False]
        ),
        tuple(
            ["AutoShotV1 (Gaussian)"]
            + [f1_text(comparison_f1("autoshot_v1_gaussian", dataset)) for dataset in ("shot", "bbc", "clipshots")]
            + [False]
        ),
        tuple(
            ["AutoShotV2"]
            + [f1_text(comparison_f1("autoshotv2_deploy", dataset)) for dataset in ("shot", "bbc", "clipshots")]
            + [True]
        ),
    ]
    col_widths = [Inches(3.7), Inches(1.8), Inches(1.8), Inches(1.8)]
    start_x = Inches(0.6)
    row_h = Inches(0.55)
    table_top = Inches(2.0)

    x = start_x
    add_rect(slide, x, table_top, sum(col_widths, Emu(0)), row_h,
             fill=COLOR_ACCENT)
    for w, h in zip(col_widths, headers):
        add_text(slide, x + Inches(0.15), table_top + Inches(0.12),
                 w - Inches(0.3), row_h,
                 h, size=13, bold=True, color=COLOR_WHITE, font=FONT_TITLE)
        x += w

    for i, row in enumerate(rows):
        name = row[0]
        vals = row[1:4]
        hi = row[4]
        x = start_x
        y = table_top + row_h * (i + 1)
        bg = RGBColor(0xE7, 0xF0, 0xFF) if hi else (COLOR_PANEL if i % 2 == 0 else COLOR_WHITE)
        add_rect(slide, x, y, sum(col_widths, Emu(0)), row_h, fill=bg)
        add_text(slide, x + Inches(0.15), y + Inches(0.12),
                 col_widths[0] - Inches(0.3), row_h,
                 name, size=13, bold=hi, color=COLOR_INK)
        xx = x + col_widths[0]
        for j, v in enumerate(vals):
            color = COLOR_ACCENT if hi else COLOR_INK
            add_text(slide, xx + Inches(0.15), y + Inches(0.12),
                     col_widths[j + 1] - Inches(0.3), row_h,
                     v, size=13, bold=hi, color=color, font=FONT_TITLE)
            xx += col_widths[j + 1]

    # Analysis
    add_text(slide, Inches(0.6), Inches(5.3), Inches(12.1), Inches(0.4),
             "Phân tích", size=13, bold=True, color=COLOR_ACCENT, font=FONT_TITLE)
    add_bullets(slide, Inches(0.6), Inches(5.75), Inches(12.1), Inches(1.6), [
        ("Mạnh trên SHOT (miền mục tiêu): ", f"F1 tốt nhất ({f1_text(RESULTS['summary']['shot_f1'])})."),
        (
            "Vượt baseline trên BBC: ",
            f"AutoShotV2 ({f1_text(RESULTS['summary']['bbc_f1'])}) cao hơn "
            f"TransNetV2 ({f1_text(comparison_f1('transnetv2_reported', 'bbc'))}) và "
            f"AutoShot tự chạy ({f1_text(comparison_f1('autoshot_reproduced_legacy', 'bbc'))}).",
        ),
        ("Giảm trên ClipShots (−3.4%): ", "post-processing tối ưu cho SHOT chưa chuyển dịch sang miền video web đa dạng."),
    ], size=13, line_spacing=1.3)


def slide_ablation(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Ablation & minh họa dự đoán", "Chương 4", 24)

    # Left: ablation table
    headers = ["Cấu hình", "F1 SHOT", "Δ vs AutoShot"]
    control_f1 = ablation_f1("A1_phase2_bce_onehot", "shot")
    rows = [
        ("A1 — BCE + one-hot", f1_text(control_f1), "—"),
        (
            "B4 — Temperature + Gaussian",
            f1_text(ablation_f1("B4_temperature_gaussian", "shot")),
            f"+{(ablation_f1('B4_temperature_gaussian', 'shot') - control_f1) * 100:.1f}%",
        ),
        (
            "B5 — Full candidate",
            f1_text(ablation_f1("B5_full_candidate", "shot")),
            f"+{(ablation_f1('B5_full_candidate', 'shot') - control_f1) * 100:.1f}%",
        ),
    ]
    col_widths = [Inches(3.6), Inches(1.3), Inches(1.5)]
    start_x = Inches(0.6)
    row_h = Inches(0.55)
    table_top = Inches(2.0)

    x = start_x
    add_rect(slide, x, table_top, sum(col_widths, Emu(0)), row_h,
             fill=COLOR_ACCENT)
    for w, h in zip(col_widths, headers):
        add_text(slide, x + Inches(0.12), table_top + Inches(0.12),
                 w - Inches(0.24), row_h,
                 h, size=12, bold=True, color=COLOR_WHITE, font=FONT_TITLE)
        x += w

    for i, row in enumerate(rows):
        x = start_x
        y = table_top + row_h * (i + 1)
        hi = (i == len(rows) - 1)
        bg = RGBColor(0xE7, 0xF0, 0xFF) if hi else (COLOR_PANEL if i % 2 == 0 else COLOR_WHITE)
        add_rect(slide, x, y, sum(col_widths, Emu(0)), row_h, fill=bg)
        for j, (w, val) in enumerate(zip(col_widths, row)):
            color = COLOR_ACCENT if (hi and j > 0) else COLOR_INK
            bold = (j == 0) or hi
            add_text(slide, x + Inches(0.12), y + Inches(0.12),
                     w - Inches(0.24), row_h,
                     val, size=12, bold=bold, color=color,
                     font=FONT_TITLE if j > 0 else FONT_MAIN)
            x += w

    # Notes under table
    add_text(slide, Inches(0.6), Inches(4.0), Inches(6.4), Inches(0.4),
             "Quan sát", size=12, bold=True,
             color=COLOR_ACCENT, font=FONT_TITLE)
    add_bullets(slide, Inches(0.6), Inches(4.45), Inches(6.4), Inches(2.8), [
        ("Gaussian smoothing: ", "cải tiến ổn định nhất qua các miền."),
        ("Focal + many-hot + temp: ", "tăng rõ trên SHOT."),
        ("Trade-off: ", "cấu hình tối ưu cho SHOT làm giảm ClipShots."),
    ], size=12, line_spacing=1.3)

    # Right: 2 example thumbnails (using transition images)
    add_text(slide, Inches(7.5), Inches(2.0), Inches(5.2), Inches(0.4),
             "Ví dụ boundary dự đoán", size=12, bold=True,
             color=COLOR_ACCENT, font=FONT_TITLE)
    examples = [
        ("cut_transition.png", "Cut — AutoShotV2 dự đoán đúng (p ≈ 0.93)."),
        ("dissolve.png",       "Dissolve — nhiều khung cùng có p cao; Gaussian giữ đỉnh."),
    ]
    ex_top = Inches(2.45)
    for i, (fname, caption) in enumerate(examples):
        y = ex_top + Inches(2.4) * i
        add_round_rect(slide, Inches(7.5), y, Inches(5.2), Inches(2.15),
                       fill=COLOR_PANEL)
        p = os.path.join(IMG, fname)
        if os.path.exists(p):
            add_image_fit(slide, p, Inches(7.65), y + Inches(0.1),
                          Inches(3.0), Inches(1.95))
        add_text(slide, Inches(10.7), y + Inches(0.4),
                 Inches(2.1), Inches(1.4),
                 caption, size=11, color=COLOR_INK)


def slide_conclusion(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Kết luận & đóng góp", "Chương 5", 25)

    # Top: 3 summary cards
    cards = [
        ("Tổng quan hệ thống",
         "Khảo sát SBD từ phương pháp cổ điển đến NAS."),
        ("Pipeline AutoShotV2",
         "Focal + many-hot + temp scaling + Gaussian trên backbone NAS đóng băng."),
        ("Hiệu năng vượt trội",
         "F1 = 0.855 trên SHOT — SOTA cho short-form."),
    ]
    card_w = Inches(3.9)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    top = Inches(2.0)
    for i, (name, desc) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        add_round_rect(slide, x, top, card_w, Inches(2.6),
                       fill=COLOR_PANEL)
        add_text(slide, x + Inches(0.3), top + Inches(0.3),
                 Inches(0.9), Inches(0.8),
                 f"0{i+1}", size=40, bold=True,
                 color=COLOR_ACCENT, font=FONT_TITLE)
        add_text(slide, x + Inches(0.3), top + Inches(1.15),
                 card_w - Inches(0.6), Inches(0.5),
                 name, size=16, bold=True, color=COLOR_INK, font=FONT_TITLE)
        add_text(slide, x + Inches(0.3), top + Inches(1.6),
                 card_w - Inches(0.6), Inches(1.0),
                 desc, size=12, color=COLOR_SUB)

    # Bottom: key message
    add_rect(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.7),
             fill=COLOR_INK)
    add_text(slide, Inches(0.85), Inches(5.2), Inches(11.7), Inches(0.4),
             "Thông điệp chính", size=11, bold=True,
             color=COLOR_ACCENT2, font=FONT_TITLE)
    add_text(slide, Inches(0.85), Inches(5.6), Inches(11.7), Inches(1.0),
             "Không cần thay đổi backbone — chỉ với classification head + "
             "hiệu chỉnh xác suất + hậu xử lý phù hợp, ta đã cải thiện F1 "
             "trên video ngắn một cách rõ rệt.",
             size=16, color=COLOR_WHITE)


def slide_limits(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Hạn chế", "Chương 5", 26)

    limits = [
        ("Chi phí GPU cho NAS",
         "Tìm kiếm kiến trúc vẫn tốn hàng chục giờ GPU — khó mở rộng."),
        ("False positive với flash & motion mạnh",
         "Camera shake / flash vẫn nhầm thành boundary."),
        ("Recall gradual còn thấp",
         "Chuyển cảnh dần khó hơn cut — cần đặc trưng dài hạn hơn."),
        ("Domain shift trên ClipShots",
         "Cấu hình tối ưu cho SHOT giảm F1 trên ClipShots −3.4pp; recall gradual chỉ 4.4% là nút thắt chính."),
        ("Không đầy đủ ablation đa miền",
         "Chưa chạy lại ma trận ablation trên tất cả dataset do giới hạn compute."),
    ]
    top = Inches(2.0)
    row_h = Inches(0.95)
    for i, (name, desc) in enumerate(limits):
        y = top + row_h * i
        add_round_rect(slide, Inches(0.6), y, Inches(1.0), Inches(0.85),
                       fill=COLOR_ACCENT2)
        add_text(slide, Inches(0.6), y + Inches(0.2),
                 Inches(1.0), Inches(0.5),
                 f"0{i+1}", size=22, bold=True,
                 color=COLOR_WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
        add_text(slide, Inches(1.8), y + Inches(0.05),
                 Inches(10.9), Inches(0.4),
                 name, size=16, bold=True, color=COLOR_INK, font=FONT_TITLE)
        add_text(slide, Inches(1.8), y + Inches(0.45),
                 Inches(10.9), Inches(0.45),
                 desc, size=12, color=COLOR_SUB)


def slide_future(prs):
    slide = blank_slide(prs)
    add_chrome(slide, "Hướng phát triển", "Chương 5", 27)

    directions = [
        ("Real-time", "Prune • Quantize • Distill cho edge device."),
        ("Multimodal",
         "Tích hợp audio + motion vector tăng độ bền vững."),
        ("Domain adaptation",
         "Thích ứng đa nền tảng (TikTok ↔ YouTube ↔ phim dài)."),
        ("Fine-grained",
         "Nhận dạng loại gradual cụ thể (Dissolve vs Wipe)."),
        ("Self-supervised",
         "Tận dụng video chưa gán nhãn — pretext boundary-aware."),
    ]
    cols = 3
    rows_n = 2
    cell_w = Inches(4.0)
    cell_h = Inches(2.35)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    top = Inches(2.0)

    for i, (name, desc) in enumerate(directions):
        c = i % cols
        r = i // cols
        x = start_x + (cell_w + gap) * c
        y = top + (cell_h + gap) * r
        add_round_rect(slide, x, y, cell_w, cell_h, fill=COLOR_PANEL)
        add_rect(slide, x, y, Inches(0.2), cell_h, fill=COLOR_ACCENT)
        add_text(slide, x + Inches(0.4), y + Inches(0.3),
                 cell_w - Inches(0.6), Inches(0.5),
                 name, size=18, bold=True, color=COLOR_INK, font=FONT_TITLE)
        add_line(slide, x + Inches(0.4), y + Inches(0.9),
                 x + Inches(1.5), y + Inches(0.9),
                 color=COLOR_ACCENT, width=2.0)
        add_text(slide, x + Inches(0.4), y + Inches(1.1),
                 cell_w - Inches(0.6), cell_h - Inches(1.2),
                 desc, size=13, color=COLOR_SUB)


def slide_thanks(prs):
    slide = blank_slide(prs)
    # Full colored background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_ACCENT)
    add_rect(slide, Inches(10.8), Inches(5.8), Inches(1.2), Inches(1.2),
             fill=COLOR_ACCENT2)
    add_rect(slide, Inches(12.1), Inches(6.5), Inches(0.5), Inches(0.5),
             fill=COLOR_WHITE)

    add_text(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.4),
             "Q & A", size=14, bold=True, color=COLOR_ACCENT2,
             font=FONT_TITLE)

    add_text(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.0),
             "Cảm ơn Hội đồng\nđã lắng nghe.",
             size=72, bold=True, color=COLOR_WHITE, font=FONT_TITLE)

    add_line(slide, Inches(0.85), Inches(5.0), Inches(3.0), Inches(5.0),
             color=COLOR_WHITE, width=2.5)

    add_text(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.4),
             "Sẵn sàng nhận câu hỏi từ Hội đồng.",
             size=18, color=COLOR_WHITE)

    add_text(slide, Inches(0.8), Inches(6.2), Inches(8.0), Inches(0.35),
             "Tặng Sềnh Mành — 22120202   •   Đỗ Tiến Mạnh — 22120203",
             size=12, color=COLOR_WHITE)
    add_text(slide, Inches(0.8), Inches(6.6), Inches(8.0), Inches(0.35),
             "GVHD: TS. Trần Thái Sơn   •   HCMUS  •  04/2026",
             size=12, color=COLOR_WHITE)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build(output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,               # 1
        slide_toc,                 # 2
        slide_summary,             # 3
        slide_motivation,          # 4
        slide_problem,             # 5
        slide_transitions,         # 6
        slide_challenges,          # 7
        slide_rq_contrib,          # 8
        slide_sbd_history,         # 9
        slide_baselines,           # 10
        slide_gaps,                # 11
        slide_arch_overview,       # 12
        slide_nas,                 # 13
        slide_loss,                # 14
        slide_calibration,         # 15
        slide_training,            # 16
        slide_diff,                # 17
        slide_dataset_shot,        # 18
        slide_dataset_challenges,  # 19
        slide_dataset_others,      # 20
        slide_metrics,             # 21
        slide_results_main,        # 22
        slide_results_cross,       # 23
        slide_ablation,            # 24
        slide_conclusion,          # 25
        slide_limits,              # 26
        slide_future,              # 27
    ]
    assert len(builders) == TOTAL_SLIDES
    for fn in builders:
        fn(prs)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    print(f"[OK] Saved: {output_path}")
    print(f"[INFO] Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        default=os.path.join(HERE, "..", "build", "AutoShotV2_Defense.pptx"),
    )
    build(parser.parse_args().output)
